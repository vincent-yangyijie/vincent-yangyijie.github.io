from pptx import Presentation
import os

def replace_font(pptx_file):
    prs = Presentation(pptx_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = '宋体'
    prs.save(pptx_file)

if __name__ == "__main__":
    folder_path = '.'  # 可修改为实际文件夹路径
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            if file.endswith('.pptx'):
                file_path = os.path.join(root, file)
                replace_font(file_path)